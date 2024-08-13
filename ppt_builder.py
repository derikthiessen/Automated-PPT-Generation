from sqlalchemy import create_engine
import pandas as pd
from datetime import datetime, timedelta
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os
from dotenv import load_dotenv

from SQL_Queries.charting_queries import queries as charting_queries
from Slide_Builders.Single_Portfolio_AUM_Slide import Single_Portfolio_AUM_Slide as spAUM_slide
from Slide_Builders.Total_Fund_AUM_Slide import Total_Fund_AUM_Slide as tfAUM_slide
from Slide_Builders.All_Stock_Returns_Slide import All_Stock_Returns_Slide as asr_slide
from Slide_Builders.Strategy_Comparison_Slide import Strategy_Comparison_Slide as sc_slide

load_dotenv()

class PPT:
    COMPOUNDING_PERIODS = 52 # Assume weekly compounding periods
    FUND_FIRST_YEAR = 1991 # derik_trading_company was established in 1991
    NUM_SLIDES = 5 # Presentations of this format will have five slides
    
    @classmethod
    def find_first_friday(cls, year):
        first_day = datetime(year, 1, 1)
        first_day_type = first_day.weekday()
        days_until_friday = (4 - first_day_type) % 7
        days_until_friday = days_until_friday + 7 if days_until_friday < 0 else days_until_friday
        return first_day + timedelta(days = days_until_friday)
    
    @classmethod
    def initialize_class_dates(cls): 
        cls.FUND_FIRST_TRADING_DAY = cls.find_first_friday(cls.FUND_FIRST_YEAR)
        cls.FUND_LAST_TRADING_DAY = datetime(2024, 5, 31)
        cls.FUND_TRADING_WEEKS = int((cls.FUND_LAST_TRADING_DAY - cls.FUND_FIRST_TRADING_DAY).days / 7)

    def __init__(self,  
                 ppt_directory: str,
                 image_directory: str,
                 authentication: str,
                 benchmark_rate: float = 0.08,
                 queries: dict = None
                 ):
        
        PPT.initialize_class_dates()

        self.ppt_directory = ppt_directory
        self.image_directory = image_directory

        self.benchmark_rate = benchmark_rate
        self.benchmark_return_sequence = self.get_benchmark_series(self.benchmark_rate)

        self.total_fund_AUM_query = None # Line chart of the entire portfolio's AUM since inception
        self.single_portfolio_AUM_query = None # Line chart for a single portfolio's AUM since its inception
        self.strategy_comparison_query = None # Double bar chart of the strategies and their average performances
        self.all_stock_returns_query = None # Scatter plot of each stock's annualized return vs when it IPO'd
        self.queries = self.preprocess_queries(queries)

        self.total_fund_AUM_data = self.load_from_postgresql(self.total_fund_AUM_query, authentication)
        self.single_portfolio_AUM_data = self.load_from_postgresql(self.single_portfolio_AUM_query, authentication)
        self.strategy_comparison_data = self.load_from_postgresql(self.strategy_comparison_query, authentication)
        self.all_stock_returns_data = self.load_from_postgresql(self.all_stock_returns_query, authentication)

        self.single_portfolio_AUM_slide = spAUM_slide(self.single_portfolio_AUM_data, 
                                                      benchmark_rate, 
                                                      self.benchmark_return_sequence, 
                                                      self.image_directory)
        self.total_fund_AUM_slide = tfAUM_slide(self.total_fund_AUM_data,
                                                self.image_directory,
                                                self.benchmark_rate
                                                )
        self.all_stock_returns_slide = asr_slide(self.all_stock_returns_data,
                                                 self.image_directory)
        
        self.strategy_comparison_slide = sc_slide(self.strategy_comparison_data,
                                                  self.image_directory)

        self.slide_titles = [self.single_portfolio_AUM_slide.title,
                             self.single_portfolio_AUM_slide.title,
                             self.total_fund_AUM_slide.title,
                             self.all_stock_returns_slide.title,
                             self.strategy_comparison_slide.title]
        self.ppt = self.build_ppt()
        print(f'Presentation complete, find it at {self.ppt}')
    
    def preprocess_queries(self, queries):
        if not queries:
            return
        
        for key, value in queries.items():
            if hasattr(self, key):
                setattr(self, key, value)

        return queries

    def load_from_postgresql(self, query, authentication):
        if not query:
            return None
        
        engine = create_engine(authentication)
        data = pd.read_sql(query, engine)
        return data
    
    def get_benchmark_series(self, benchmark_rate):
        weekly_return_rate = (1 + benchmark_rate) ** (1 / PPT.COMPOUNDING_PERIODS) - 1
        benchmark_series = [(1 + weekly_return_rate) ** i for i in range(1, PPT.FUND_TRADING_WEEKS + 1)]
        return benchmark_series

    def build_ppt(self):
        prs = Presentation()  

        num_slides = PPT.NUM_SLIDES

        for i in range(num_slides):
            prs = self.apply_default_slide_format(prs, self.slide_titles[i])

        for i in range(num_slides):
            prs = self.build_slide(i, prs)

        print('Beginning to save the new presentation')
        ppt_file = 'Derik_Trading_Company_Report.pptx'
        path = os.path.join(self.ppt_directory, ppt_file)
        prs.save(path)
        print('Finished saving the new presentation')

        return path
    
    def apply_default_slide_format(self, prs: Presentation, title: str) -> Presentation:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        # Remove the title placeholder
        for shape in slide.shapes:
            if shape.is_placeholder:
                sp = shape.placeholder_format.idx
                if sp == 0:  
                    sp = shape
                    slide.shapes._spTree.remove(sp._element)

        # Add a blue rectangle across the top of the slide
        top_rectangle = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,  # Shape type
            Cm(0),                # Left position
            Cm(0),                # Top position
            prs.slide_width,      # Width 
            Cm(2)                 # Height 
        )
        # Set the fill color of the rectangle to blue
        fill = top_rectangle.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 89, 179)

        # Add text inside the rectangle
        text_frame = top_rectangle.text_frame
        text_frame.text = title

        # Adjust text properties
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(28)  # Set font size to 28 points (adjust as needed)
                run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white for contrast
                run.font.bold = True

            # Center-align the text
            paragraph.alignment = PP_ALIGN.CENTER

        # Center-align the text frame within the rectangle
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Add logo image to the bottom left corner
        slide.shapes.add_picture(os.path.join(self.image_directory, 'RBC_Logo.png'), 
                                Cm(0.2),  # Left position
                                Cm(17.62), # Top position
                                Cm(0.98), # Width
                                Cm(1.28) # Height
                                )

        return prs
    
    def build_slide(self, slide_num: int, prs: Presentation) -> Presentation:
        if slide_num == 0:
            prs = self.build_slide_1(prs)
        elif slide_num == 1:
            prs = self.build_slide_2(prs)
        elif slide_num == 2:
            prs = self.build_slide_3(prs)
        elif slide_num == 3:
            prs = self.build_slide_4(prs)
        else:
            prs = self.build_slide_5(prs)
        
        return prs

    def build_slide_1(self, prs: Presentation) -> Presentation:
        slide = prs.slides[0]

        # Add line chart image
        slide.shapes.add_picture(self.single_portfolio_AUM_slide.line_chart_path, 
                                 Cm(2.3), # Left position
                                 Cm(2.54), # Top position
                                 Cm(20.8), # Width
                                 Cm(14.85) # Height
                                )
        
        return prs
    
    def build_slide_2(self, prs: Presentation) -> Presentation:
        
        slide = prs.slides[1]

        # Add bar chart image
        slide.shapes.add_picture(self.single_portfolio_AUM_slide.bar_chart_path,
                                 Cm(0.7), # Left position
                                 Cm(4.7),  # Top position
                                 Cm(16.06), # Width
                                 Cm(10.71) # Height
                                 )
        
        benchmark_rate = self.single_portfolio_AUM_slide.benchmark_rate * 100
        annualized_rate = self.single_portfolio_AUM_slide.annualized_portfolio_return * 100
        year = self.single_portfolio_AUM_slide.year_established
        num_stocks = self.single_portfolio_AUM_slide.num_stocks_held
        start_capital = self.single_portfolio_AUM_slide.starting_capital

        text = [f'Benchmark return: {benchmark_rate}%',
                f'Annualized portfolio return: {annualized_rate:,.2f}%',
                f'Year established: {year}',
                f'Stocks currently held: {num_stocks}',
                f'Starting capital: ${start_capital:,.0f}']
        
        text = '\n\n'.join(text)

        textbox = slide.shapes.add_textbox(Cm(15.2), # Left position
                                           Cm(5.14), # Top position
                                           Cm(9.5), # Width
                                           Cm(9.5) # Height
                                          )
        
        textframe = textbox.text_frame
        textframe.word_wrap = True

        p = textframe.add_paragraph()
        p.text = text
        p.font.size = Pt(18)

        return prs
        
    def build_slide_3(self, prs: Presentation) -> Presentation:
        slide = prs.slides[2]

        # Add line chart image
        slide.shapes.add_picture(self.total_fund_AUM_slide.line_chart_path, 
                                 Cm(2.3), # Left position
                                 Cm(2.54), # Top position
                                 Cm(20.8), # Width
                                 Cm(14.85) # Height
                                )
        
        return prs
    
    def build_slide_4(self, prs: Presentation) -> Presentation:
        slide = prs.slides[3]

        # Add scatter chart image
        slide.shapes.add_picture(self.all_stock_returns_slide.scatter_chart_path,
                                  Cm(2.3), # Left position
                                  Cm(2.54), # Top position
                                  Cm(20.8), # Width
                                  Cm(14.85) # Height
                                  )
        
        return prs
    
    def build_slide_5(self, prs: Presentation) -> Presentation:
        slide = prs.slides[4]

        slide.shapes.add_picture(self.strategy_comparison_slide.bar_chart_path,
                                 Cm(2.3), # Left position
                                 Cm(2.54), # Top position
                                 Cm(20.8), # Width
                                 Cm(14.85) # Height
                                 )
        
        return prs

authentication = os.getenv('DATABASE_URL')
ppt_directory = os.getenv('PPT_DIRECTORY')
image_directory = os.getenv('IMAGE_DIRECTORY')

test = PPT(ppt_directory = ppt_directory, image_directory = image_directory, authentication = authentication, queries = charting_queries)